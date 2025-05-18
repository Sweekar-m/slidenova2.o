import sys
import os
from pptx import Presentation
from google import genai
from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QLineEdit, QPushButton, QVBoxLayout,
    QHBoxLayout, QFormLayout, QSpinBox, QLabel, QFileDialog, QComboBox, QCheckBox
)
from PyQt6.QtGui import QPixmap, QFont
from PyQt6.QtCore import Qt

folder_path = ""
temp = ""
topic = ""
slidenum = ""
save = ""
name = ""
ex = ""
def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller .exe """
    try:
        base_path = sys._MEIPASS  # Temporary folder created by PyInstaller
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)
class SlideNova(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Slide Nova - Professional Presentation Generator")
        self.setMinimumSize(1200, 800)
        self.use_dark_mode = False  # Default theme is light
        self.dropdowns = []  # To store dropdown widgets for templates
        self.dropdown_sets = [
        [resource_path(f"temp1/temp{i}.png") for i in range(1, 4)],
        [resource_path(f"temp2/temp2_{i}.png") for i in range(1, 4)],
        [resource_path(f"temp3/temp3_{i}.jpg") for i in range(1, 4)],
        [resource_path(f"temp4/temp4_{i}.jpg") for i in range(1, 4)]  
        ]

        # Main Widget and Layout
        main_widget = QWidget()
        self.setCentralWidget(main_widget)
        self.main_layout = QVBoxLayout(main_widget)
        self.main_layout.setContentsMargins(20, 20, 20, 20)

        # Header Section
        self.setup_header()

        # Input and Preview Section
        self.setup_input_preview_section()

        # Template Dropdown Section
        self.setup_template_dropdown_section()

        # Generate Button
        self.setup_generate_button()

        # Theme Switch
        self.setup_theme_switch()

        # Footer Section
        self.setup_footer()

        # Apply Default Theme
        self.apply_theme()

    def setup_header(self):
        header_layout = QVBoxLayout()
        header_layout.setAlignment(Qt.AlignmentFlag.AlignCenter)

        heading = QLabel("Slide Nova")
        heading.setFont(QFont("Arial", 48, QFont.Weight.Bold))
        heading.setAlignment(Qt.AlignmentFlag.AlignCenter)

        subheading = QLabel("Create Stunning Presentations Effortlessly")
        subheading.setFont(QFont("Arial", 16, QFont.Weight.Normal))
        subheading.setAlignment(Qt.AlignmentFlag.AlignCenter)

        header_layout.addWidget(heading)
        header_layout.addWidget(subheading)
        self.main_layout.addLayout(header_layout)
        self.main_layout.addSpacing(20)

    def setup_input_preview_section(self):
        input_preview_layout = QHBoxLayout()
        input_preview_layout.setSpacing(40)

        # Input Section
        input_layout = QFormLayout()
        input_layout.setFormAlignment(Qt.AlignmentFlag.AlignCenter)
        input_layout.setVerticalSpacing(20)

        topic_label = QLabel("Enter Topic:")
        topic_label.setFont(QFont("Arial", 18))
        self.topic_input = QLineEdit()
        self.topic_input.setPlaceholderText("Enter the topic for your presentation")
        self.topic_input.setFont(QFont("Arial", 16))
        self.topic_input.setFixedWidth(400)
        input_layout.addRow(topic_label, self.topic_input)

        slide_label = QLabel("Number of Slides:")
        slide_label.setFont(QFont("Arial", 18))
        self.slide_spin = QSpinBox()
        self.slide_spin.setRange(1, 50)
        self.slide_spin.setFont(QFont("Arial", 16))
        self.slide_spin.setFixedWidth(100)
        input_layout.addRow(slide_label, self.slide_spin)

        input_preview_layout.addLayout(input_layout)

        # Preview Section
        preview_layout = QVBoxLayout()
        preview_layout.setAlignment(Qt.AlignmentFlag.AlignCenter)

        preview_label_heading = QLabel("Template Preview")
        preview_label_heading.setFont(QFont("Arial", 18, QFont.Weight.Bold))
        preview_label_heading.setAlignment(Qt.AlignmentFlag.AlignCenter)

        self.preview_label = QLabel()
        self.preview_label.setFixedSize(500, 250)
        self.preview_label.setStyleSheet("border: 2px solid #ddd; background-color: #f9f9f9; border-radius: 10px;")
        self.preview_label.setAlignment(Qt.AlignmentFlag.AlignCenter)

        preview_layout.addWidget(preview_label_heading)
        preview_layout.addWidget(self.preview_label)
        preview_layout.setContentsMargins(0, -20, 0, 0)  # Move the preview section slightly upward
        input_preview_layout.addLayout(preview_layout)

        self.main_layout.addLayout(input_preview_layout)
        self.main_layout.addSpacing(20)

    def setup_template_dropdown_section(self):
        dropdown_layout = QHBoxLayout()
        dropdown_layout.setSpacing(20)
        dropdown_layout.setAlignment(Qt.AlignmentFlag.AlignCenter)

        category_names = ["Gradient", "Dark & Edgy", "Elegant Swirl", "Modern Abstract"]
        self.template_names = [
            ["Deep Twilight", "Deep Charcoal", "Sunset Glow"],
            ["Deep Haze", "Cyber Nexus", "Deep Sync"],
            ["Aqua Flow", "Crystal Edge", "Polygonal Ice"],
            ["Business Glow", "Digital Frame", "Sleek Design"],
        ]

        for category_name, image_set, template_names in zip(category_names, self.dropdown_sets, self.template_names):
            vbox = QVBoxLayout()
            label = QLabel(category_name)
            label.setFont(QFont("Arial", 14, QFont.Weight.Bold))
            label.setAlignment(Qt.AlignmentFlag.AlignCenter)

            dropdown = QComboBox()
            for name, img in zip(template_names, image_set):
                dropdown.addItem(name, img)  # Associate template name with its image path

            dropdown.setFixedWidth(200)
            dropdown.currentIndexChanged.connect(self.handle_dropdown_change)  # Connect to handler
            self.dropdowns.append(dropdown)

            vbox.addWidget(label)
            vbox.addWidget(dropdown)
            dropdown_layout.addLayout(vbox)

        self.main_layout.addLayout(dropdown_layout)
        self.main_layout.addSpacing(20)

    def handle_dropdown_change(self):
        """
        Handles the dropdown value change. Fetches and prints the selected value and image path.
        """
        global temp
        for dropdown in self.dropdowns:
            if self.sender() == dropdown:
                temp = dropdown.currentText()  # Get the selected template name
                
                self.update_image_preview(dropdown)
                break

    def update_image_preview(self, dropdown):
        """
        Updates the image preview based on the selected template in the dropdown.
        """
        image_path = dropdown.currentData()
        if image_path and os.path.exists(image_path):
            pixmap = QPixmap(image_path).scaled(
                self.preview_label.width(),
                self.preview_label.height(),
                Qt.AspectRatioMode.KeepAspectRatio,
            )
            self.preview_label.setPixmap(pixmap)
        else:
            self.preview_label.clear()

    def setup_generate_button(self):
        self.generate_button = QPushButton("Generate Presentation")
        self.generate_button.setFont(QFont("Arial", 16, QFont.Weight.Bold))
        self.generate_button.setStyleSheet("""
            QPushButton {
                background-color: #2F4F4F;
                color: white;
                border-radius: 20px;
                padding: 10px 20px;
            }
            QPushButton:hover {
                background-color: #3E606F;
            }
        """)
        self.generate_button.clicked.connect(self.generate_presentation)

        self.main_layout.addWidget(self.generate_button, alignment=Qt.AlignmentFlag.AlignCenter)
        self.main_layout.addSpacing(20)

    def setup_theme_switch(self):
        theme_switch_layout = QHBoxLayout()
        theme_switch_layout.setAlignment(Qt.AlignmentFlag.AlignRight)

        theme_label = QLabel("Dark Mode:")
        theme_label.setFont(QFont("Arial", 14))
        theme_switch_layout.addWidget(theme_label)

        self.theme_switch = QCheckBox()
        self.theme_switch.setStyleSheet("""
            QCheckBox::indicator {
                width: 40px;
                height: 20px;
                border-radius: 10px;
                background-color: #ccc;
            }
            QCheckBox::indicator:checked {
                background-color: #2F4F4F;
            }
        """)
        self.theme_switch.stateChanged.connect(self.toggle_theme)
        theme_switch_layout.addWidget(self.theme_switch)

        self.main_layout.addLayout(theme_switch_layout)

    def setup_footer(self):
        footer = QLabel('Developed by <a href="https://sweekar-m-portfolio.vercel.app/">Sweekar M</a>')
        footer.setFont(QFont("Arial", 12))
        footer.setAlignment(Qt.AlignmentFlag.AlignCenter)
        footer.setOpenExternalLinks(True)
        self.main_layout.addWidget(footer, alignment=Qt.AlignmentFlag.AlignBottom)

    def generate_presentation(self):
        global temp,topic,slidenum,folder_path
        topic = self.topic_input.text()
        slidenum = self.slide_spin.value()
        file_dialog = QFileDialog()
        folder_path = file_dialog.getExistingDirectory(self, "Select Folder")
        print("The selected template is:", temp)
        print(topic)
        print(slidenum)
        if folder_path:
            print(f"Selected folder: {folder_path}")
            generated_info=self.code_generate()
            generated_info=generated_info.replace("```python", "")
            generated_info=generated_info.replace("```", "")
            print(generated_info)
            
            code1= '''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller .exe """
    try:
        base_path = sys._MEIPASS  # Temporary folder created by PyInstaller
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)
# Function to set the title slide properly
def set_title_slide(slide, title_text, subtitle_text):
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title_text
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].font.size = Pt(48)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = [rgb text]
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_shape.text = subtitle_text
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.paragraphs[0].font.size = Pt(28)
    subtitle_tf.paragraphs[0].font.color.rgb = [rgb text]
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Function to set the title of a slide
def set_title(slide, title_text, font_size):
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    text_frame = title_shape.text_frame
    text_frame.text = title_text
    
    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = [rgb text]
    p.alignment = PP_ALIGN.CENTER

# Function to set body content with bullet points
def set_body_content(slide, body_text_list, font_size):
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for item in body_text_list:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = [rgb text]
        p.alignment = PP_ALIGN.LEFT

# Function to set the background with a gradient color
def set_background(slide, start_color, end_color):
    fill = slide.background.fill
    fill.gradient()
    stops = fill.gradient_stops
    stop1 = stops[0]
    stop1.position = 0.0
    stop1.color.rgb = start_color
    stop2 = stops[1]
    stop2.position = 1.0
    stop2.color.rgb = end_color

# Function to create a presentation
def create_presentation(topic_details):
    prs = Presentation()
    prs.slide_width = Inches(10)   # Width for 16:9 ratio
    prs.slide_height = Inches(5.625)  

    # Slide 1 (Title slide)
    slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    set_title_slide(slide_1, topic_details['title'], topic_details['subtitle'])
    set_background(slide_1,[rgb value])  # Dark blue to light blue gradient
    
    # Loop over topic details and create content slides
    for slide_detail in topic_details['content_slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        set_title(slide, slide_detail['title'], 45)
        set_body_content(slide, slide_detail['content'], 25)
        set_background(slide,[rgb value])  # Dark gray to light gray gradient

    # Save the presentation
    prs.save(r'[path]')
    print("Presentation created successfully!")

# Define topic details as a dictionary
[value]

# Create the presentation for the given topic
create_presentation(topic_details)

            '''  
            code2='''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller .exe """
    try:
        base_path = sys._MEIPASS  # Temporary folder created by PyInstaller
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)
# Function to set the background image in the foreground
def set_foreground_image(slide, image_path, prs):
    img = slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    slide.shapes._spTree.remove(img._element)  # Move image to the foreground
    slide.shapes._spTree.append(img._element)

# Function to set the title and subtitle in the exact center
def set_title_slide(slide, title_text, subtitle_text):
    text_color = RGBColor(255, 255, 255) 

    # Title Formatting (Centered)
    title_shape = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(1))
    title_tf = title_shape.text_frame
    title_tf.text = title_text
    p = title_tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    # Subtitle Formatting (Centered)
    subtitle_shape = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1))
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.text = subtitle_text
    p = subtitle_tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

# Function to set the title of a slide
def set_title(slide, title_text, is_dark_bg):
    text_color = RGBColor(255, 255, 255) if is_dark_bg else RGBColor(0, 0, 0)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    text_frame = title_shape.text_frame
    text_frame.text = title_text
    
    p = text_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

# Function to set body content with bullet points
def set_body_content(slide, body_text_list, is_dark_bg):
    text_color = RGBColor(255, 255, 255) if is_dark_bg else RGBColor(0, 0, 0)

    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for item in body_text_list:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(25)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.LEFT

# Function to create the presentation
def create_presentation(topic_details, title_slide_image_path, bg_image_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Slide 1 (Title Slide with a different image)
    slide_1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Slide
    set_foreground_image(slide_1, title_slide_image_path, prs)  # Set different image for title slide
    set_title_slide(slide_1, topic_details['title'], topic_details['subtitle'])

    # Content Slides (Using the common background image)
    for slide_detail in topic_details['content_slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Slide
        set_foreground_image(slide, bg_image_path, prs)  # Image in foreground
        set_title(slide, slide_detail['title'], is_dark_bg=True)  # Fixed missing argument
        set_body_content(slide, slide_detail['content'], is_dark_bg=True)

    # Save the presentation
    prs.save("[path]")
    print("Presentation created successfully!")

# Define topic details
[topic]

# Paths for images
title_slide_image_path =r"[image_path_first]"  # Different image for first slide
bg_image_path = r"[image_path_nxt]"  # Common background for other slides

# Generate Presentation
create_presentation(topic_details, title_slide_image_path, bg_image_path)

            '''
            code3='''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller .exe """
    try:
        base_path = sys._MEIPASS  # Temporary folder created by PyInstaller
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path) 
# Function to set the background image in the foreground
def set_foreground_image(slide, image_path, prs):
    img = slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    slide.shapes._spTree.remove(img._element)  # Move image to the foreground
    slide.shapes._spTree.append(img._element)

# Function to set the title and subtitle in the exact center
def set_title_slide(slide, title_text, subtitle_text):
    text_color = RGBColor(0, 0, 0) 

    # Title Formatting (Centered)
    title_shape = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(1))
    title_tf = title_shape.text_frame
    title_tf.text = title_text
    p = title_tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    # Subtitle Formatting (Centered)
    subtitle_shape = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1))
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.text = subtitle_text
    p = subtitle_tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

# Function to set the title of a slide
def set_title(slide, title_text, is_dark_bg):
    text_color = RGBColor(0, 0, 0) if is_dark_bg else RGBColor(0, 0, 0)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    text_frame = title_shape.text_frame
    text_frame.text = title_text
    
    p = text_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

# Function to set body content with bullet points
def set_body_content(slide, body_text_list, is_dark_bg):
    text_color = RGBColor(0, 0, 0) if is_dark_bg else RGBColor(0, 0, 0)

    textbox = slide.shapes.add_textbox(Inches(1.5), Inches(0.9), Inches(7.5), Inches(5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for item in body_text_list:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.LEFT

# Function to create the presentation
def create_presentation(topic_details, bg_image_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Slide 1 (Title Slide with the same background image)
    slide_1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Slide
    set_foreground_image(slide_1, bg_image_path, prs)  # Use the same image for all slides
    set_title_slide(slide_1, topic_details['title'], topic_details['subtitle'])

    # Content Slides (Using the same background image)
    for slide_detail in topic_details['content_slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Slide
        set_foreground_image(slide, bg_image_path, prs)  # Use the same image for all slides
        set_title(slide, slide_detail['title'], is_dark_bg=True)
        set_body_content(slide, slide_detail['content'], is_dark_bg=True)

    # Save the presentation
    prs.save("[path]")
    print("Presentation created successfully!")

# Define topic details
[topic]

# Path for the common background image
bg_image_path = r"[image_path]"

# Generate Presentation
create_presentation(topic_details, bg_image_path)

            '''
            code4='''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller .exe """
    try:
        base_path = sys._MEIPASS  # Temporary folder created by PyInstaller
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)
# Function to set the background image in the foreground
def set_foreground_image(slide, image_path, prs):
    img = slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    slide.shapes._spTree.remove(img._element)  # Move image to the foreground
    slide.shapes._spTree.append(img._element)

# Function to set the title and subtitle in the exact center
def set_title_slide(slide, title_text, subtitle_text):
    text_color = RGBColor(255, 255, 255) 

    # Title Formatting (Centered)
    title_shape = slide.shapes.add_textbox(Inches(3), Inches(1.8), Inches(4), Inches(1))
    title_tf = title_shape.text_frame
    title_tf.text = title_text
    p = title_tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER 

    # Subtitle Formatting (Centered)
    subtitle_shape = slide.shapes.add_textbox(Inches(3), Inches(3.2), Inches(4), Inches(1))
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.text = subtitle_text
    p = subtitle_tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

# Function to set the title of a slide
def set_title(slide, title_text, is_dark_bg):
    text_color = RGBColor(0, 0, 0) if is_dark_bg else RGBColor(0, 0, 0)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    text_frame = title_shape.text_frame
    text_frame.text = title_text
    
    p = text_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

# Function to set body content with bullet points
def set_body_content(slide, body_text_list, is_dark_bg):
    text_color = RGBColor(0, 0, 0) if is_dark_bg else RGBColor(0, 0, 0)

    textbox = slide.shapes.add_textbox(Inches(2.8), Inches(1), Inches(7.12), Inches(5.13))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for item in body_text_list:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(25)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.LEFT

# Function to create the presentation
def create_presentation(topic_details, title_slide_image_path, bg_image_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Slide 1 (Title Slide with a different image)
    slide_1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Slide
    set_foreground_image(slide_1, title_slide_image_path, prs)  # Set different image for title slide
    set_title_slide(slide_1, topic_details['title'], topic_details['subtitle'])

    # Content Slides (Using the common background image)
    for slide_detail in topic_details['content_slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Slide
        set_foreground_image(slide, bg_image_path, prs)  # Image in foreground
        set_title(slide, slide_detail['title'], is_dark_bg=True)  # Fixed missing argument
        set_body_content(slide, slide_detail['content'], is_dark_bg=True)

    # Save the presentation
    prs.save("[path]")
    print("Presentation created successfully!")

# Define topic details
[topic]

# Paths for images
title_slide_image_path =r"[image_path_first]"  # Different image for first slide
bg_image_path = r"[image_path_nxt]"  # Common background for other slides

# Generate Presentation
create_presentation(topic_details, title_slide_image_path, bg_image_path)
'''
        
            code5='''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller .exe """
    try:
        base_path = sys._MEIPASS  # Temporary folder created by PyInstaller
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

# Function to set the background image in the foreground
def set_foreground_image(slide, image_path, prs):
    img = slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    slide.shapes._spTree.remove(img._element)  # Move image to the foreground
    slide.shapes._spTree.append(img._element)

# Function to set the title and subtitle in the exact center
def set_title_slide(slide, title_text, subtitle_text):
    text_color = RGBColor(255, 255, 255) 

    # Title Formatting (Centered)
    title_shape = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(1))
    title_tf = title_shape.text_frame
    title_tf.text = title_text
    p = title_tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    # Subtitle Formatting (Centered)
    subtitle_shape = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1))
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.text = subtitle_text
    p = subtitle_tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

# Function to set the title of a slide
def set_title(slide, title_text, is_dark_bg):
    text_color = RGBColor(255, 255, 255) if is_dark_bg else RGBColor(0, 0, 0)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    text_frame = title_shape.text_frame
    text_frame.text = title_text
    
    p = text_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

# Function to set body content with bullet points
def set_body_content(slide, body_text_list, is_dark_bg):
    text_color = RGBColor(255, 255, 255) if is_dark_bg else RGBColor(0, 0, 0)

    textbox = slide.shapes.add_textbox(Inches(1.5), Inches(0.9), Inches(7.5), Inches(5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for item in body_text_list:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.LEFT

# Function to create the presentation
def create_presentation(topic_details, bg_image_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Slide 1 (Title Slide with the same background image)
    slide_1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Slide
    set_foreground_image(slide_1, bg_image_path, prs)  # Use the same image for all slides
    set_title_slide(slide_1, topic_details['title'], topic_details['subtitle'])

    # Content Slides (Using the same background image)
    for slide_detail in topic_details['content_slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Slide
        set_foreground_image(slide, bg_image_path, prs)  # Use the same image for all slides
        set_title(slide, slide_detail['title'], is_dark_bg=True)
        set_body_content(slide, slide_detail['content'], is_dark_bg=True)

    # Save the presentation
    prs.save("[path]")
    print("Presentation created successfully!")

# Define topic details
[topic]

# Path for the common background image
bg_image_path = r"[image_path]"

# Generate Presentation
create_presentation(topic_details, bg_image_path)

'''
            code1=code1.replace("[topic]",generated_info)
            code2=code2.replace("[topic]",generated_info)
            code3=code3.replace("[topic]",generated_info)
            code4=code4.replace("[topic]",generated_info)
            code5=code5.replace("[topic]",generated_info)
            name=topic.split()
            ex=name[0]
            if(folder_path=="D:/" or folder_path=="C:/" ):
                global save
              
                print(ex)
                save=folder_path+ex+".pptx"
                code1=code1.replace("[path]", save)
                code2=code2.replace("[path]", save)
                code3=code3.replace("[path]", save)
                code4=code4.replace("[path]", save)
                code5=code5.replace("[path]", save)
                print(save)
            else:
                save=folder_path+"/"+ex+".pptx"
                print(ex)
                code1=code1.replace("[path]", save)
                code2=code2.replace("[path]", save)
                code3=code3.replace("[path]", save)
                code4=code4.replace("[path]", save)
                code5=code5.replace("[path]", save)
                print(save)

            if temp == "Deep Twilight":
                code1=code1.eplace("[rgb value]","RGBColor(10, 10, 20), RGBColor(30, 30, 50)")
                code1=code1.replace("[rgb text]","RGBColor(250, 250, 250)")
                print(code1)
                # self.code_execute(code1)

            elif temp == "Deep Charcoal":
                code1=code1.replace("[rgb value]","RGBColor(45, 45, 45), RGBColor(75, 75, 75)")
                code1=code1.replace("[rgb text]","RGBColor(250, 250, 250)")
                print(code1)
                self.code_execute(code1)
            elif temp == "Sunset Glow":
                code1=code1.replace("[rgb value]","RGBColor(255, 94, 98), RGBColor(255, 195, 113)")
                code1=code1.replace("[rgb text]","RGBColor(0, 4, 0)")
                print(code1)
                self.code_execute(code1)
            else:
                if temp == "Deep Haze":
                    code2 = code2.replace("[image_path_first]",resource_path(f"temp2/temp2_1.png") )
                    code2 = code2.replace("[image_path_nxt]",resource_path(f"temp2/nxt.png"))
                    print(code2)
                    self.code_execute(code2)


                elif temp == "Cyber Nexus":
                    code5 = code5.replace("[image_path]", resource_path(f"temp2/temp2_2.png") )
                    print(code5)
                    self.code_execute(code5)

                elif temp=="Deep Sync":
                    code2 = code2.replace("[image_path_first]", resource_path(f"temp2/temp2_3.png") )
                    code2 = code2.replace("[image_path_nxt]", resource_path(f"temp2/nex2.png"))
                    print(code2)
                    self.code_execute(code2)

                elif temp == "Polygonal Ice":
                    code3 = code3.replace("[image_path]", resource_path(f"temp3/temp3_3.jpg"))
                    print(code3)
                    self.code_execute(code3)

                elif temp == "Aqua Flow":
                    code3 = code3.replace("[image_path]", resource_path(f"temp3/temp3_1.jpg"))
                    print(code3)
                    self.code_execute(code3)
                elif temp == "Crystal Edge":
                    code3 = code3.replace("[image_path]", resource_path(f"temp3/temp3_2.jpg"))
                    print(code3)
                    self.code_execute(code3)
                elif temp == "Business Glow":
                    code3 = code3.replace("[image_path]", resource_path(f"temp4/temp4_1.jpg"))
                    print(code3)
                    self.code_execute(code3)
                elif temp == "Digital Fram":
                    code4 = code4.replace("[image_path_first]", resource_path(f"temp4/temp4_2.jpg"))
                    code4 = code4.replace("[image_path_nxt]", resource_path(f"temp4/nxt.jpg"))
                    print(code4)
                    self.code_execute(code4)
                elif temp == "Sleek Design":
                    code3 = code3.replace("[image_path]", resource_path(f"temp4/temp4_3.jpg"))
                    print(code3)
                    self.code_execute(code3)





    def code_execute(self, script):
        namespace = {}
        try:
            # Execute the script in the provided namespace
            print(script)
            exec(script, namespace)
            
            # Check if the 'create_presentation' function exists in the namespace
            if 'create_presentation' in namespace and 'topic_details' in namespace:
                print("Executing 'create_presentation' with 'topic_details'...")
                try:
                    namespace['create_presentation'](namespace['topic_details'])
                except Exception as function_error:
                    print(f"Error while calling 'create_presentation': {function_error}")
            else:
                print("Error: 'create_presentation' or 'topic_details' not found in the script namespace.")

            # Open the file if it was successfully created
            if os.path.exists(save):
                os.startfile(save)
            else:
                print(f"Error: File not found at path '{save}'.")

        except SyntaxError as syntax_error:
            print(f"Syntax Error in the script: {syntax_error}")
        except Exception as e:
            print(f"Error executing the script: {e}")

    def code_generate(self):
        global folder_path,topic,slidenum
        prompt='''
           Give me the information on the topic [topic] in the same format up to [number] slide(s). Each slide can have paragraph or point-wise information. If point-wise, include 4 to 5 points. The last slide should contain the conclusion. Give the result in the exact format as shown below, without any explanations.

If the content is in paragraph form, enclose it in [" "] brackets. Do not include explanations or formatting like **bold**. Output the result as a Python dictionary with the variable name `topic_details` and in the format given below.

Example format:

topic_details = {
    "title": "Chat GPT",
    "subtitle": "Understanding and Utilizing the Power of AI Conversations",
    "content_slides": [
        {"title": "What is Chat GPT?", "content": [
            "* A cutting-edge language model developed by OpenAI.",
            "* Designed for engaging in natural and interactive conversations.",
            "* Leverages deep learning techniques for advanced text generation.",
            "* Capable of understanding context and providing relevant responses.",
            "* Continuously learns and adapts based on user interactions.",
            "* Can be used for various applications, from customer service to content creation."
        ]},
        {"title": "The Technology Behind Chat GPT", "content": [
            "* Built upon the GPT (Generative Pre-trained Transformer) architecture.",
            "* Utilizes a massive dataset of text and code for pre-training.",
            "* Employs a transformer network with self-attention mechanisms.",
            "* Learns patterns and relationships within the data to predict the next word in a sequence.",
            "* Fine-tuned through supervised learning and reinforcement learning to enhance performance.",
            "* Its ability to understand and generate human-like text is constantly evolving."
        ]}
    ]
}


    '''
        prompt = prompt.replace("[topic]", topic)
        prompt = prompt.replace("[path]", folder_path)
        prompt = prompt.replace("[number]", str(slidenum))

        
        client = genai.Client(api_key="AIzaSyAQAOHfnwmnB5qvVCsOrKGuEWc93ISh78o")
        response = client.models.generate_content(
            model="gemini-2.0-flash", contents=prompt
        )
        return response.text

    def toggle_theme(self):
        self.use_dark_mode = self.theme_switch.isChecked()
        self.apply_theme()

    def apply_theme(self):
        if self.use_dark_mode:
            self.setStyleSheet("""
                QWidget {
                    background: qlineargradient(
                        x1: 0, y1: 0, x2: 0, y2: 1,
                        stop: 0 #2F4F4F, stop: 1 #000000
                    );
                    color: white;
                }
                QLabel {
                    color: white;
                }
            """)
        else:
            self.setStyleSheet("""
                QWidget {
                    background: qlineargradient(
                        x1: 0, y1: 0, x2: 0, y2: 1,
                        stop: 0 #f9f9f9, stop: 1 #ffffff
                    );
                    color: black;
                }
                QLabel {
                    color: black;
                }
            """)


if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = SlideNova()
    window.show()
    sys.exit(app.exec())